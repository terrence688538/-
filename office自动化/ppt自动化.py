from pptx import Presentation                          #新建ppt
from pptx.util import Inches,Pt

ppt=Presentation()                                     #建立一个空的幻灯片
slide=ppt.slides.add_slide(ppt.slide_layouts[1])       #在ppt插入一个幻灯片
body_shape=slide.shapes.placeholders                   #获取到所有的文本框
body_shape[0].text='这是占位符[0]'
body_shape[1].text='这是占位符[1]'
new_paragraph=body_shape[1].text_frame.add_paragraph()
new_paragraph.text='新段落'
new_paragraph.font.bold=True
new_paragraph.font.italic=True                          #字体倾斜
new_paragraph.font.size=Pt(15)                          #字体大小
new_paragraph.font.underline=True                       #文字下划线

#添加一个新的文本框
left=Inches(2)
top=Inches(2)
width=Inches(3)
height=Inches(3)
textbox=slide.shapes.add_textbox(left,top,width,height)
textbox.text='这是新文本框'
new_para=textbox.text_frame.add_paragraph()
new_para.text='这是文本框里的第二段'

#添加图片
pptx=Presentation()
slide=pptx.slides.add_slide(pptx.slide_layouts[1])
left=Inches(1)
top=Inches(1)
width=Inches(2)
height=Inches(2)
pic=slide.shapes.add_picture('C:/Users/Mechrevo/Desktop/7.png',left,top,width,height)

#添加表格
rows=2
cols=2
left=Inches(2)
top=Inches(2)
width=Inches(4)
height=Inches(4)
table=slide.shapes.add_table(rows,cols,left,top,width,height).table         #表格的大小
table.columns[0].width=Inches(1)                                            #列宽
table.columns[0].width=Inches(3)
table.cell(0,0).text='1'
table.cell(0,1).text='2'
table.cell(1,0).text='3'
table.cell(1,1).text='4'

ppt.save('C:/Users/Mechrevo/Desktop/test.pptx')
pptx.save('C:/Users/Mechrevo/Desktop/test.pptx')

#修改已有的ppt

pptv=Presentation('C:/Users/Mechrevo/Desktop/model.pptx')
shape=pptv.slides[0].shapes                                     #获取这一页所有的文本框
for i in shape:
    print(i.text)
    print('\n')
print(shape[1].text_frame.paragraphs[1].text)
print(shape[1].text_frame.paragraphs[1].runs[0].text)
shape[1].text_frame.paragraphs[1].runs[0].text='更改后的新段落'
pptv.save('C:/Users/Mechrevo/Desktop/test_0.pptx')
